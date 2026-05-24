import os
import zipfile
import xml.etree.ElementTree as ET
import re
from pptx import Presentation
from pypdf import PdfReader

def extract_pptx_text(file_path: str):
    """
    Extracts text slide-by-slide from a PPTX file, including text boxes, tables, and speaker notes.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    prs = Presentation(file_path)
    slide_contents = []
    
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        
        # 1. Extract text from standard shapes (text frames)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_run = "".join(run.text for run in paragraph.runs if run.text)
                    if text_run.strip():
                        slide_text.append(text_run.strip())
            
            # 2. Extract text from tables if present
            if shape.has_table:
                table_text = []
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_cells.append(cell.text.strip())
                    if row_cells:
                        table_text.append(" | ".join(row_cells))
                if table_text:
                    slide_text.append("\n".join(table_text))
        
        # 3. Extract speaker notes (helps ground AI questions with lecture details)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"\n[Ghi chú slide: {notes}]")
                
        # Join texts and save slide representation
        combined_text = "\n".join(slide_text).strip()
        slide_contents.append({
            "slide_num": idx + 1,
            "content": combined_text
        })
        
    return slide_contents

def extract_pdf_text(file_path: str):
    """
    Extracts text page-by-page from a PDF file.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    reader = PdfReader(file_path)
    slide_contents = []
    
    for idx, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        slide_contents.append({
            "slide_num": idx + 1,
            "content": text.strip()
        })
        
    return slide_contents

def extract_slide_content(file_path: str):
    """
    Wrapper function to automatically detect file type and extract slide content.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        return extract_pptx_text(file_path)
    elif ext == ".pdf":
        return extract_pdf_text(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}. Only .pptx and .pdf are supported.")

def is_red_or_highlighted(r_el, namespaces):
    rPr = r_el.find('w:rPr', namespaces)
    if rPr is None:
        return False
        
    color_el = rPr.find('w:color', namespaces)
    if color_el is not None:
        val = color_el.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')
        if val:
            val = val.upper()
            if val.startswith(("FF", "EE", "DD")) or val == "RED":
                return True
                
    highlight_el = rPr.find('w:highlight', namespaces)
    if highlight_el is not None:
        val = highlight_el.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')
        if val and val.lower() in ("red", "lightred"):
            return True
            
    return False

def extract_docx_raw_text(file_path: str) -> str:
    """
    Extracts raw text paragraph-by-paragraph from a DOCX file, wrapping colored/highlighted
    runs in <correct>...</correct> tags to serve as context for the AI parser.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    with zipfile.ZipFile(file_path) as z:
        xml_content = z.read('word/document.xml')
        root = ET.fromstring(xml_content)
        
        namespaces = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        }
        
        paragraphs = root.findall('.//w:p', namespaces)
        
        lines = []
        for p in paragraphs:
            text_parts = []
            for r in p.findall('.//w:r', namespaces):
                t = r.find('w:t', namespaces)
                if t is not None and t.text:
                    if is_red_or_highlighted(r, namespaces):
                        text_parts.append(f"<correct>{t.text}</correct>")
                    else:
                        text_parts.append(t.text)
            p_text = "".join(text_parts).strip()
            if p_text:
                lines.append(p_text)
                
        return "\n".join(lines)

def split_paragraph_content(p_text: str):
    # 1. First split by any embedded "Câu [số]"
    sub_parts = re.split(r'(?=\b(?:C\u00e2u|C\u00c2U|c\u00e2u)\s*\d+[\.:\s])', p_text)
    
    final_parts = []
    for part in sub_parts:
        part = part.strip()
        if not part:
            continue
            
        # 2. Find all matches of option prefixes, possibly preceded by <correct> tag
        pattern = r'(?:<correct>\s*)?\b([A-Da-d])[\.\)]'
        matches = list(re.finditer(pattern, part))
        
        pos_to_split = []
        for m in matches:
            start_idx = m.start()
            letter = m.group(1).upper()
            
            # Avoid matching names like C. Mác
            snippet = part[start_idx:start_idx+20]
            snippet_clean = re.sub(r'</?correct>', '', snippet)
            if letter == 'C' and re.match(r'^C\.\s*M\u00e1c', snippet_clean, re.IGNORECASE):
                continue
            if re.match(r'^[A-Da-d][\.\)]\s*M\u00e1c', snippet_clean, re.IGNORECASE):
                continue
                
            pos_to_split.append(start_idx)
            
        if not pos_to_split:
            final_parts.append(part)
        else:
            last_pos = 0
            for pos in pos_to_split:
                chunk = part[last_pos:pos].strip()
                if chunk:
                    final_parts.append(chunk)
                last_pos = pos
            chunk = part[last_pos:].strip()
            if chunk:
                final_parts.append(chunk)
                
    return final_parts

def extract_docx_quiz(file_path: str):
    """
    Parses a DOCX file directly and extracts multiple choice questions,
    options, and correct answers (detected by red font or highlight).
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    with zipfile.ZipFile(file_path) as z:
        xml_content = z.read('word/document.xml')
        root = ET.fromstring(xml_content)
        
        namespaces = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        }
        
        paragraphs = root.findall('.//w:p', namespaces)
        
        doc_paragraphs = []
        for p in paragraphs:
            text_parts = []
            for r in p.findall('.//w:r', namespaces):
                text_el = r.find('w:t', namespaces)
                if text_el is not None and text_el.text:
                    txt = text_el.text
                    if is_red_or_highlighted(r, namespaces):
                        text_parts.append(f"<correct>{txt}</correct>")
                    else:
                        text_parts.append(txt)
            
            p_text = "".join(text_parts).strip()
            if p_text:
                # Preprocess inline options with missing spaces
                # e.g., "giai đoạnB. 2 giai đoạn" -> "giai đoạn B. 2 giai đoạn"
                p_text = re.sub(r'([^\s\.\>])([A-Da-d][\.\)])', r'\1 \2', p_text)
                
                # Split paragraph content if it contains inline questions or options
                split_parts = split_paragraph_content(p_text)
                for part in split_parts:
                    doc_paragraphs.append(part)
                
        raw_questions = []
        current_q = None
        
        def is_question_start(text):
            return re.match(r'^(?:[Cc]âu|CÂU)\s*\d+', text, re.IGNORECASE) is not None
            
        def get_option_start_prefix(text):
            clean = re.sub(r'</?correct>', '', text).strip()
            m = re.match(r'^([A-Da-d])[\.\)\s/]', clean)
            if m:
                return m.group(1).upper()
            return None

        # Parse paragraphs
        for p in doc_paragraphs:
            if is_question_start(p):
                if current_q:
                    raw_questions.append(current_q)
                current_q = {
                    "text": p,
                    "options": []
                }
            elif current_q and get_option_start_prefix(p):
                current_q["options"].append(p)
            else:
                if current_q:
                    if len(current_q["options"]) > 0:
                        current_q["options"][-1] += "\n" + p
                    else:
                        current_q["text"] += "\n" + p
                        
        if current_q:
            raw_questions.append(current_q)
            
        structured_questions = []
        for rq in raw_questions:
            q_text = rq["text"]
            options_raw = rq["options"]
            
            q_text_clean = re.sub(r'</?correct>', '', q_text).strip()
            
            options = []
            correct_answer = None
            
            should_parse_inline = False
            if len(options_raw) < 4:
                full_block = q_text + "\n" + "\n".join(options_raw)
                clean_block = re.sub(r'</?correct>', '', full_block)
                prefixes_found = 0
                for pref in ['A', 'B', 'C', 'D']:
                    if re.search(r'\b' + pref + r'[\.\)]', clean_block):
                        prefixes_found += 1
                if prefixes_found >= 3:
                    should_parse_inline = True
            
            if should_parse_inline:
                full_block = q_text + "\n" + "\n".join(options_raw)
                clean_block = re.sub(r'</?correct>', '', full_block)
                
                inline_matches = []
                last_pos = 0
                for pref in ['A', 'B', 'C', 'D']:
                    pattern = r'\b' + pref + r'[\.\)\s/]'
                    for m in re.finditer(pattern, clean_block[last_pos:]):
                        pos = last_pos + m.start()
                        matched_text = clean_block[pos:pos+10]
                        if pref == 'C' and re.match(r'^C\.\s*Mác', matched_text, re.IGNORECASE):
                            # Skip C. Mác name
                            continue
                        inline_matches.append((pref, pos))
                        last_pos = pos + len(m.group(0))
                        break
                
                opt_pattern = r'(<correct>)?\s*\b([A-Da-d])([\.\)\s/])\s*(</correct>)?'
                orig_matches = list(re.finditer(opt_pattern, full_block))
                
                filtered_matches = []
                next_expected = 'A'
                for m in orig_matches:
                    pref = m.group(2).upper()
                    if pref == next_expected:
                        start_pos = m.start()
                        snippet = re.sub(r'</?correct>', '', full_block[start_pos:start_pos+15])
                        if pref == 'C' and re.match(r'^C\.\s*Mác', snippet, re.IGNORECASE):
                            continue
                        filtered_matches.append(m)
                        if next_expected == 'A': next_expected = 'B'
                        elif next_expected == 'B': next_expected = 'C'
                        elif next_expected == 'C': next_expected = 'D'
                        elif next_expected == 'D': next_expected = 'DONE'
                
                if len(filtered_matches) > 0:
                    q_text_clean = re.sub(r'</?correct>', '', full_block[:filtered_matches[0].start()]).strip()
                    
                    for i, match in enumerate(filtered_matches):
                        start = match.start()
                        end = filtered_matches[i+1].start() if i + 1 < len(filtered_matches) else len(full_block)
                        opt_str = full_block[start:end].strip()
                        
                        is_correct = "<correct>" in opt_str
                        clean_opt = re.sub(r'</?correct>', '', opt_str).strip()
                        clean_opt = re.sub(r'^[A-Da-d][\.\)\s/]\s*', '', clean_opt).strip()
                        
                        if clean_opt:
                            options.append(clean_opt)
                            if is_correct:
                                correct_answer = clean_opt
            else:
                # Multi-paragraph options (Style 1)
                for opt_str in options_raw:
                    is_correct = "<correct>" in opt_str
                    clean_opt = re.sub(r'</?correct>', '', opt_str).strip()
                    clean_opt = re.sub(r'^[A-Da-d][\.\)\s/]\s*', '', clean_opt).strip()
                    if clean_opt:
                        options.append(clean_opt)
                        if is_correct:
                            correct_answer = clean_opt
                            
            if len(options) > 0:
                if not correct_answer:
                    correct_answer = options[0]
                    
                q_display = re.sub(r'^(?:[Cc]âu|CÂU)\s*\d+[\.:\s]*', '', q_text_clean).strip()
                
                structured_questions.append({
                    "question": q_display,
                    "options": options,
                    "correct_answer": correct_answer,
                    "explanation": "Được nhập tự động từ file Word.",
                    "source_slide": 1
                })
                
        return structured_questions

